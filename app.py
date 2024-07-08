# from flask import Flask, request, render_template, send_file
# from pptx import Presentation
# from pptx.util import Pt
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from io import BytesIO

# app = Flask(__name__)

# def replace_text(shape, placeholder, new_text, font_name="Segoe UI", font_size=8):
#     for paragraph in shape.text_frame.paragraphs:
#         for run in paragraph.runs:
#             if placeholder in run.text:
#                 run.text = run.text.replace(placeholder, new_text)
#                 run.font.name = font_name
#                 run.font.size = Pt(font_size)

# def update_student_info(shape, students_info):
#     text_frame = shape.text_frame
#     text_frame.clear()

#     for i, student in enumerate(students_info):
#         p = text_frame.add_paragraph()
#         if i == 0:
#             p.text = f": {student}"
#         else:
#             p.text = f"  {student}"
#         p.space_after = Pt(6)  # Adjust the space after each paragraph for smaller spacing
#         p.font.name = "Segoe UI"
#         p.font.size = Pt(8)  # Adjust font size as needed


# def update_application_list(shape, applications):
#     text_frame = shape.text_frame
#     text_frame.clear()

#     for application in applications:
#         p = text_frame.add_paragraph()
#         p.text = f"• {application}"
#         p.font.name = "Segoe UI"
#         p.font.size = Pt(8)

# @app.route('/', methods=['GET', 'POST'])
# def index():
#     if request.method == 'POST':
#         # Get the form data
#         user_program = request.form['program']
#         user_project = request.form['project']
#         user_domain = request.form['domain']
#         user_project_info = request.form['project_info']
#         user_guide = request.form['guide']
#         num_students = int(request.form['num_students'])
#         students_info = [request.form[f'roll_no_{i+1}'] + " - " + request.form[f'student_name_{i+1}'] for i in range(num_students)]
#         num_applications = int(request.form['num_applications'])
#         applications = [request.form[f'application_name_{i+1}'] for i in range(num_applications)]
#         num_screenshots = int(request.form['num_screenshots'])

#         # Load the presentation
#         presentation = Presentation('template.pptx')

#         # Initialize counters and store content of text fields
#         text_fields_content = []
#         for slide in presentation.slides:
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     text_fields_content.append((shape, shape.text))

#         # Update the text fields with the user's input
#         for shape, content in text_fields_content:
#             replace_text(shape, "Electronics and Telecommunication Engineering", user_program, font_size=14)
#             replace_text(shape, "Project Name", user_project)
#             replace_text(shape, "Domain Name", user_domain)
#             replace_text(shape, "Guide-Name", user_guide)
#             replace_text(shape, "Describe-the-project-briefly", user_project_info)

#             if "Roll No. - Student Name_1" in content:
#                 update_student_info(shape, students_info)

#             if "Application 1" in content:
#                 update_application_list(shape, applications)

#         # Replace the placeholder image with the specified image
#         if 'image' in request.files:
#             image_file = request.files['image']
#             placeholder_name = "Picture 11"
#             image_stream = BytesIO(image_file.read())
#             for slide in presentation.slides:
#                 for shape in slide.shapes:
#                     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == placeholder_name:
#                         # Remove the old picture
#                         spTree = shape._element.getparent()
#                         spTree.remove(shape._element)
#                         # Add the new picture
#                         slide.shapes.add_picture(image_stream, shape.left, shape.top, width=shape.width, height=shape.height)

#         # Replace the screenshots and remove unused placeholders
#         screenshot_placeholders = ["Picture 8", "Picture 10", "Picture 3"]
#         for i in range(3):
#             screenshot_key = f'screenshot_{i+1}'
#             placeholder_name = screenshot_placeholders[i]
#             for slide in presentation.slides:
#                 for shape in slide.shapes:
#                     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == placeholder_name:
#                         if i < num_screenshots and screenshot_key in request.files:
#                             screenshot_file = request.files[screenshot_key]
#                             screenshot_stream = BytesIO(screenshot_file.read())
#                             # Remove the old picture
#                             spTree = shape._element.getparent()
#                             spTree.remove(shape._element)
#                             # Add the new picture
#                             slide.shapes.add_picture(screenshot_stream, shape.left, shape.top, width=shape.width, height=shape.height)
#                         elif i >= num_screenshots:
#                             # Remove the unused placeholder
#                             spTree = shape._element.getparent()
#                             spTree.remove(shape._element)

#         # Save the updated presentation to a BytesIO object
#         updated_presentation = BytesIO()
#         presentation.save(updated_presentation)
#         updated_presentation.seek(0)

#         return send_file(updated_presentation, as_attachment=True, download_name='updated_presentation.pptx')

#     return render_template('index.html')

# if __name__ == '__main__':
#     app.run(debug=True)



from flask import Flask, request, render_template, send_file
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import os

app = Flask(__name__)

def replace_text(shape, placeholder, new_text, font_name="Segoe UI", font_size=8):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if placeholder in run.text:
                run.text = run.text.replace(placeholder, new_text)
                run.font.name = font_name
                run.font.size = Pt(font_size)

def update_student_info(shape, students_info):
    text_frame = shape.text_frame
    text_frame.clear()

    for i, student in enumerate(students_info):
        p = text_frame.add_paragraph()
        if i == 0:
            p.text = f": {student}"
        else:
            p.text = f"  {student}"
        p.space_after = Pt(6)  # Adjust the space after each paragraph for smaller spacing
        p.font.name = "Segoe UI"
        p.font.size = Pt(8)  # Adjust font size as needed

def update_application_list(shape, applications):
    text_frame = shape.text_frame
    text_frame.clear()

    for application in applications:
        p = text_frame.add_paragraph()
        p.text = f"• {application}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(8)

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Get the form data
        user_program = request.form['program']
        user_project = request.form['project']
        user_domain = request.form['domain']
        user_project_info = request.form['project_info']
        user_guide = request.form['guide']
        num_students = int(request.form['num_students'])
        students_info = [request.form[f'roll_no_{i+1}'] + " - " + request.form[f'student_name_{i+1}'] for i in range(num_students)]
        num_applications = int(request.form['num_applications'])
        applications = [request.form[f'application_name_{i+1}'] for i in range(num_applications)]
        num_screenshots = int(request.form['num_screenshots'])

        # Load the presentation
        presentation = Presentation('template.pptx')

        # Initialize counters and store content of text fields
        text_fields_content = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_fields_content.append((shape, shape.text))

        # Update the text fields with the user's input
        for shape, content in text_fields_content:
            replace_text(shape, "Electronics and Telecommunication Engineering", user_program, font_size=14)
            replace_text(shape, "Project Name", user_project)
            replace_text(shape, "Domain Name", user_domain)
            replace_text(shape, "Guide-Name", user_guide)
            replace_text(shape, "Describe-the-project-briefly", user_project_info)

            if "Roll No. - Student Name_1" in content:
                update_student_info(shape, students_info)

            if "Application 1" in content:
                update_application_list(shape, applications)

        # Replace the placeholder image with the specified image
        if 'image' in request.files:
            image_file = request.files['image']
            placeholder_name = "Picture 11"
            image_stream = BytesIO(image_file.read())
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == placeholder_name:
                        # Remove the old picture
                        spTree = shape._element.getparent()
                        spTree.remove(shape._element)
                        # Add the new picture
                        slide.shapes.add_picture(image_stream, shape.left, shape.top, width=shape.width, height=shape.height)

        # Replace the screenshots and remove unused placeholders
        screenshot_placeholders = ["Picture 8", "Picture 10", "Picture 3"]
        for i in range(3):
            screenshot_key = f'screenshot_{i+1}'
            placeholder_name = screenshot_placeholders[i]
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.name == placeholder_name:
                        if i < num_screenshots and screenshot_key in request.files:
                            screenshot_file = request.files[screenshot_key]
                            screenshot_stream = BytesIO(screenshot_file.read())
                            # Remove the old picture
                            spTree = shape._element.getparent()
                            spTree.remove(shape._element)
                            # Add the new picture
                            slide.shapes.add_picture(screenshot_stream, shape.left, shape.top, width=shape.width, height=shape.height)
                        elif i >= num_screenshots:
                            # Remove the unused placeholder
                            spTree = shape._element.getparent()
                            spTree.remove(shape._element)

        # Save the updated presentation to a BytesIO object
        updated_presentation = BytesIO()
        presentation.save(updated_presentation)
        updated_presentation.seek(0)

        return send_file(updated_presentation, as_attachment=True, download_name='updated_presentation.pptx')

    return render_template('index.html')

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
